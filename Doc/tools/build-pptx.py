#!/usr/bin/env python3
"""
Ensambla Doc/Fauchard_Presentacion_Comercial.pptx a partir de los PNG generados
por render-slides.mjs. Una slide 16:9 por imagen, ajustada (contain) y centrada
sobre el fondo de marca (#020617), para verse idéntico al HTML.
"""
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

HERE = Path(__file__).resolve().parent
SLIDES = HERE / "slides"
OUT = HERE.parent / "Fauchard_Presentacion_Comercial.pptx"
BG = RGBColor(0x02, 0x06, 0x17)

# Lienzo 16:9 widescreen (13.333" x 7.5") en EMU.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def main():
    manifest = json.loads((SLIDES / "manifest.json").read_text())

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # layout en blanco

    for name in manifest:
        img_path = SLIDES / name
        with Image.open(img_path) as im:
            iw, ih = im.size

        slide = prs.slides.add_slide(blank)

        # Fondo de marca sólido.
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

        # Contain: escalar para que la imagen entre completa, centrada.
        scale = min(SLIDE_W / iw, SLIDE_H / ih)
        w = int(iw * scale)
        h = int(ih * scale)
        left = int((SLIDE_W - w) / 2)
        top = int((SLIDE_H - h) / 2)
        slide.shapes.add_picture(str(img_path), Emu(left), Emu(top), Emu(w), Emu(h))

    prs.save(str(OUT))
    print(f"OK · {len(manifest)} slides -> {OUT}")


if __name__ == "__main__":
    main()
