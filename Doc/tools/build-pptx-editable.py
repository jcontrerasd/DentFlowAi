#!/usr/bin/env python3
"""
Versión EDITABLE del deck: reconstruye Fauchard_Presentacion_Comercial.html como
PPTX nativo (texto y formas editables, fiel a la marca). El diagrama de flujo y
el mockup van como imagen (capturados en slides/diagram.png y slides/mockup.png).

Parsea el HTML para no re-transcribir contenido.
Salida: Doc/Fauchard_Presentacion_Comercial_Editable.pptx
"""
from pathlib import Path

from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
SLIDES = HERE / "slides"
HTML = HERE.parent / "Fauchard_Presentacion_Comercial.html"
OUT = HERE.parent / "Fauchard_Presentacion_Comercial_Editable.pptx"

# ---- Paleta de marca ----
BG = RGBColor(0x02, 0x06, 0x17)
PANEL = RGBColor(0x0F, 0x17, 0x2A)
PANEL2 = RGBColor(0x11, 0x1C, 0x33)
LINE = RGBColor(0x2A, 0x38, 0x52)
TEALSOFT = RGBColor(0x0C, 0x2B, 0x29)
FG = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
MUTED2 = RGBColor(0xCB, 0xD5, 0xE1)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
VERDE = RGBColor(0x22, 0xC5, 0x5E)
AMARILLO = RGBColor(0xEA, 0xB3, 0x08)
NARANJA = RGBColor(0xF9, 0x73, 0x16)
ROJO = RGBColor(0xEF, 0x44, 0x44)

BODY = "Inter"
TITLE = "Georgia"

SW, SH = 13.333, 7.5
MX = 0.7  # margen lateral

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def new_slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = BG
    return s


def _set_run(r, text, size, color, bold, font, italic):
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color


def text(slide, l, t, w, h, runs, size=14, color=FG, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.12):
    """runs: str | [(txt, opts)]  ó  lista de párrafos [[(txt,opts),...], ...]."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(2)

    if isinstance(runs, str):
        paras = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs

    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for s, o in para:
            _set_run(p.add_run(), s, o.get("size", size), o.get("color", color),
                     o.get("bold", bold), o.get("font", font), o.get("italic", False))
    return tb


def rrect(slide, l, t, w, h, fill=PANEL, line=LINE, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    # radio de esquina
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def chip(slide, l, t, txt, w=0.95, h=0.42, fg=TEAL, fill=TEALSOFT, line=TEAL, size=11):
    rrect(slide, l, t, w, h, fill=fill, line=line, radius=0.5)
    text(slide, l, t, w, h, [(txt, {"bold": True, "size": size})], color=fg,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)


def pic_contain(slide, path, l, t, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    slide.shapes.add_picture(str(path), Inches(l + (w - pw) / 2),
                             Inches(t + (h - ph) / 2), Inches(pw), Inches(ph))


def eyebrow(slide, chip_txt, rest, t=0.55):
    chip(slide, MX, t, chip_txt, w=1.15)
    text(slide, MX + 1.3, t, 9, 0.42, [(rest.upper(), {"bold": True, "size": 12})],
         color=TEAL, anchor=MSO_ANCHOR.MIDDLE, wrap=False)


def section_title(slide, title, lead=None):
    text(slide, MX, 1.05, SW - 2 * MX, 0.9, title, size=33, color=FG, font=TITLE)
    if lead:
        text(slide, MX, 1.95, SW - 2 * MX, 0.7, lead, size=14, color=MUTED2)


def clean(el):
    return el.get_text(" ", strip=True) if el else ""


# ---------- parse HTML ----------
soup = BeautifulSoup(HTML.read_text(encoding="utf-8"), "html.parser")


# ---------- 1. PORTADA ----------
def slide_portada():
    s = new_slide()
    hero = soup.select_one("header.hero")
    eb = hero.select_one(".eyebrow")
    eb_chip = clean(eb.select_one(".chip"))
    eb_rest = clean(eb).replace(eb_chip, "", 1).strip()
    eyebrow(s, eb_chip, eb_rest, t=0.7)

    # h1 con highlight teal
    h1 = hero.select_one("h1")
    hl = h1.select_one(".hl")
    hl_txt = clean(hl)
    full = clean(h1)
    before, after = full, ""
    if hl_txt and hl_txt in full:
        before, after = full.split(hl_txt, 1)
    runs = [(before, {"size": 40, "font": TITLE, "color": FG})]
    if hl_txt:
        runs.append((hl_txt, {"size": 40, "font": TITLE, "color": TEAL, "italic": True}))
        runs.append((after, {"size": 40, "font": TITLE, "color": FG}))
    text(s, MX, 1.7, SW - 2 * MX, 2.3, [runs], spacing=1.05)

    p = clean(hero.select_one("p"))
    text(s, MX, 4.1, 10.5, 1.2, p, size=17, color=MUTED2)

    badges = hero.select(".hero-badge")
    bw, gap = 2.85, 0.18
    for i, b in enumerate(badges[:4]):
        bl = MX + i * (bw + gap)
        rrect(s, bl, 5.7, bw, 0.95, fill=PANEL, line=LINE, radius=0.16)
        text(s, bl + 0.15, 5.7, bw - 0.3, 0.95, clean(b), size=12.5, color=MUTED2,
             anchor=MSO_ANCHOR.MIDDLE)


# ---------- 2. RECORRIDO (diagrama imagen) ----------
def slide_flujo():
    s = new_slide()
    sec = soup.select_one("#flujo")
    eyebrow(s, "FAUCHARD", "Entregable 1 · Recorrido completo")
    section_title(s, clean(sec.select_one("h2")), clean(sec.select_one(".lead")))
    pic_contain(s, SLIDES / "diagram.png", MX, 2.5, SW - 2 * MX, 4.7)


# ---------- 3. NIVELES (tabla nativa) ----------
def slide_niveles():
    s = new_slide()
    sec = soup.select_one("#complejidad")
    eyebrow(s, "FAUCHARD", "Entregable 2 · Clasificación")
    section_title(s, clean(sec.select_one("h2")), clean(sec.select_one(".lead")))

    colors = [VERDE, AMARILLO, NARANJA, ROJO]
    rows = sec.select("table.cx tbody tr")
    headers = ["Nivel", "Criterio que lo activa", "Lectura operativa", "Liga", "Perfil técnico"]
    nrows, ncols = len(rows) + 1, len(headers)
    top, left = 2.7, MX
    tw, th = SW - 2 * MX, 4.4
    gtable = s.shapes.add_table(nrows, ncols, Inches(left), Inches(top), Inches(tw), Inches(th))
    tbl = gtable.table
    tbl.first_row = False
    for w, cw in zip(tbl.columns, [1.7, 4.0, 2.7, 1.2, 2.3]):
        w.width = Inches(cw)

    def cell(r, c, txt, color=MUTED2, bold=False, size=10, fill=PANEL):
        cl = tbl.cell(r, c)
        cl.fill.solid()
        cl.fill.fore_color.rgb = fill
        cl.margin_left = cl.margin_right = Pt(5)
        cl.margin_top = cl.margin_bottom = Pt(3)
        cl.vertical_anchor = MSO_ANCHOR.TOP
        tf = cl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r0 = p.add_run()
        _set_run(r0, txt, size, color, bold, BODY, False)
        return cl

    for c, hd in enumerate(headers):
        cell(0, c, hd.upper(), color=MUTED2, bold=True, size=9.5, fill=PANEL2)

    for ri, tr in enumerate(rows, start=1):
        tds = tr.find_all("td")
        lvl = clean(tds[0].select_one(".lvl"))
        cell(ri, 0, lvl, color=colors[ri - 1], bold=True, size=11)
        cell(ri, 1, clean(tds[1]), size=9.5)
        cell(ri, 2, clean(tds[2]), size=9.5)
        cell(ri, 3, clean(tds[3]), color=FG, bold=True, size=10)
        cell(ri, 4, clean(tds[4]), size=9.5)


# ---------- 4. EL MOTOR ----------
def slide_motor():
    s = new_slide()
    sec = soup.select_one("#motor")
    eyebrow(s, "FAUCHARD", "Entregable 3 · El motor por dentro")
    section_title(s, clean(sec.select_one("h2")), clean(sec.select_one(".lead")))

    colw = (SW - 2 * MX - 0.6) / 3
    top = 2.6
    # Entradas
    def io_col(l, title, items, title_color=MUTED):
        text(s, l, top, colw, 0.4, [(title.upper(), {"bold": True, "size": 11})],
             color=title_color, align=PP_ALIGN.CENTER)
        y = top + 0.5
        for it in items:
            b = clean(it.select_one("b"))
            span = clean(it.select_one("span"))
            rrect(s, l, y, colw, 0.92, fill=PANEL2, line=LINE, radius=0.12)
            text(s, l + 0.12, y + 0.06, colw - 0.24, 0.85,
                 [[(b, {"bold": True, "size": 12, "color": FG})], [(span, {"size": 10.5, "color": MUTED})]],
                 spacing=1.0)
            y += 1.02

    io_col(MX, "→ Entra del caso", sec.select(".bb-in .io"))

    # Núcleo
    cl = MX + colw + 0.3
    rrect(s, cl, top, colw, 4.5, fill=TEALSOFT, line=TEAL, radius=0.08)
    text(s, cl, top + 0.2, colw, 0.5, "FAUCHARD", size=22, color=FG, font=TITLE,
         align=PP_ALIGN.CENTER)
    text(s, cl, top + 0.75, colw, 0.35, [("MOTOR DE CLASIFICACIÓN Y ASIGNACIÓN",
         {"bold": True, "size": 9})], color=TEAL, align=PP_ALIGN.CENTER)
    y = top + 1.25
    for i, d in enumerate(sec.select(".core .decision"), start=1):
        ptxt = clean(d.select_one("p"))
        rrect(s, cl + 0.15, y, colw - 0.3, 0.72, fill=BG, line=LINE, radius=0.12)
        text(s, cl + 0.28, y + 0.04, colw - 0.5, 0.66,
             [[(f"{i}. ", {"bold": True, "color": TEAL, "size": 11}), (ptxt, {"size": 10.5, "color": MUTED2})]],
             spacing=1.0)
        y += 0.8

    # Salidas
    rl = MX + 2 * (colw + 0.3)
    io_col(rl, "Entrega →", sec.select(".bb-out .io"))


# ---------- 5. EN PANTALLA (mockup imagen) ----------
def slide_pantalla():
    s = new_slide()
    sec = soup.select_one("#pantalla")
    eyebrow(s, "FAUCHARD", "Entregable 4 · En pantalla")
    section_title(s, clean(sec.select_one("h2")), clean(sec.select_one(".lead")))
    pic_contain(s, SLIDES / "mockup.png", MX, 2.55, SW - 2 * MX, 4.6)


# ---------- 6. POR QUÉ IMPORTA ----------
def slide_valor():
    s = new_slide()
    sec = soup.select_one("#valor")
    eyebrow(s, "FAUCHARD", "Entregable 5 · Por qué importa")
    section_title(s, clean(sec.select_one("h2")), clean(sec.select_one(".lead")))

    cards = sec.select(".scard")
    cw = (SW - 2 * MX - 2 * 0.25) / 3
    ch = 1.85
    top = 2.7
    for i, c in enumerate(cards[:6]):
        col = i % 3
        row = i // 3
        l = MX + col * (cw + 0.25)
        t = top + row * (ch + 0.2)
        rrect(s, l, t, cw, ch, fill=PANEL, line=LINE, radius=0.1)
        text(s, l + 0.18, t + 0.14, cw - 0.36, 0.5,
             clean(c.select_one("h3")), size=14, color=FG, bold=True)
        text(s, l + 0.18, t + 0.66, cw - 0.36, 0.75,
             clean(c.select_one("p")), size=10.5, color=MUTED2, spacing=1.05)
        text(s, l + 0.18, t + 1.42, cw - 0.36, 0.35,
             clean(c.select_one(".mech")), size=9, color=MUTED, spacing=1.0)


# ---------- 7. PARÁMETROS ----------
def slide_param_intro():
    s = new_slide()
    sec = soup.select_one("#parametros")
    eyebrow(s, "FAUCHARD", "Entregable 6 · Cómo se parametriza")
    section_title(s, clean(sec.select_one("h2")), clean(sec.select_one(".lead")))


def slide_group_divider(here, tab, gtitle, gsub):
    s = new_slide()
    chip(s, MX, 0.7, "FAUCHARD", w=1.15)
    text(s, MX + 1.3, 0.7, 9, 0.42, [("PARÁMETROS", {"bold": True, "size": 12})],
         color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text(s, MX, SH / 2 - 1.2, SW - 2 * MX, 0.5,
         [(f"Admin › Fauchard › Configuración › {here}", {"size": 14})], color=MUTED, align=PP_ALIGN.CENTER)
    text(s, MX, SH / 2 - 0.6, SW - 2 * MX, 1.0, gtitle, size=34, color=FG, font=TITLE,
         align=PP_ALIGN.CENTER)
    text(s, MX, SH / 2 + 0.5, SW - 2 * MX, 0.7, gsub, size=15, color=MUTED2, align=PP_ALIGN.CENTER)
    if tab:
        chip(s, SW / 2 - 1.2, SH / 2 + 1.3, tab, w=2.4, h=0.45, size=12)


def slide_param(card, here):
    s = new_slide()
    pn = clean(card.select_one(".pn"))
    pname = clean(card.select_one(".pname"))
    pdef = clean(card.select_one(".pdef"))
    pdesc = clean(card.select_one(".pdesc"))

    # breadcrumb arriba
    text(s, MX, 0.45, SW - 2 * MX, 0.35, [(f"Parámetros · {here}", {"size": 11})], color=MUTED)
    # header: chip N° + nombre + default
    chip(s, MX, 0.85, f"N°{pn}", w=1.0, h=0.5, size=14)
    text(s, MX + 1.15, 0.85, 8.6, 0.6, pname, size=24, color=FG, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, SW - MX - 1.9, 0.88, 1.9, 0.46, fill=BG, line=LINE, radius=0.2)
    text(s, SW - MX - 1.9, 0.88, 1.9, 0.46, [(pdef, {"bold": True, "size": 13})],
         color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False, font=TITLE)
    # descripción
    text(s, MX, 1.6, SW - 2 * MX, 0.7, pdesc, size=12, color=MUTED2, spacing=1.1)

    # grid 2x2 de campos
    fields = card.select(".pf-grid .pf")
    cw = (SW - 2 * MX - 0.3) / 2
    chh = 1.05
    gy = 2.5
    for i, pf in enumerate(fields[:4]):
        col = i % 2
        row = i // 2
        l = MX + col * (cw + 0.3)
        t = gy + row * (chh + 0.18)
        rrect(s, l, t, cw, chh, fill=PANEL2, line=LINE, radius=0.1)
        lbl = clean(pf.select_one(".pf-l"))
        val = clean(pf.select_one(".pf-v"))
        text(s, l + 0.16, t + 0.08, cw - 0.32, 0.3,
             [(lbl.upper(), {"bold": True, "size": 9.5})], color=TEAL)
        text(s, l + 0.16, t + 0.38, cw - 0.32, chh - 0.46, val, size=11, color=MUTED2, spacing=1.05)

    # ejemplo
    ej = card.select_one(".pcard-ej p")
    ejt = 4.95
    rels = card.select(".pcard-rel .relchip")
    ejh = 1.55 if not rels else 1.25
    if ej:
        rrect(s, MX, ejt, SW - 2 * MX, ejh, fill=PANEL, line=LINE, radius=0.07)
        text(s, MX + 0.18, ejt + 0.1, SW - 2 * MX - 0.36, 0.3,
             [("EJEMPLO", {"bold": True, "size": 9.5})], color=TEAL)
        text(s, MX + 0.18, ejt + 0.42, SW - 2 * MX - 0.36, ejh - 0.5,
             clean(ej), size=11, color=MUTED2, spacing=1.1)
    # relacionados
    if rels:
        ry = ejt + ejh + 0.12
        text(s, MX, ry, 1.4, 0.4, [("RELACIONADOS", {"bold": True, "size": 9})],
             color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        x = MX + 1.5
        for rc in rels:
            t = clean(rc)
            w = min(3.6, 0.13 * len(t) + 0.3)
            chip(s, x, ry, t, w=w, h=0.4, size=9.5)
            x += w + 0.15


def slides_parametros():
    slide_param_intro()
    sec = soup.select_one("#parametros")
    here = ""
    for el in sec.find_all(recursive=False):
        cls = el.get("class", [])
        if "param-path" in cls:
            here = clean(el.select_one(".here"))
            tab = clean(el.select_one(".tabchip"))
            # el group-h viene justo después
            gh = el.find_next_sibling("div", class_="group-h")
            gtitle = clean(gh.select_one("h3")) if gh else here
            gsub = clean(gh.select_one("span")) if gh else ""
            slide_group_divider(here, tab, gtitle, gsub)
        elif "pstack" in cls:
            for card in el.select(".pcard"):
                slide_param(card, here)


# ---------- 8. CIERRE ----------
def slide_cierre():
    s = new_slide()
    cl = soup.select_one(".closing")
    text(s, MX, SH / 2 - 1.0, SW - 2 * MX, 1.2, clean(cl.select_one("h2")),
         size=40, color=FG, font=TITLE, align=PP_ALIGN.CENTER)
    text(s, MX, SH / 2 + 0.4, SW - 2 * MX, 1.0, clean(cl.select_one("p")),
         size=17, color=MUTED2, align=PP_ALIGN.CENTER)


# ---------- build ----------
slide_portada()
slide_flujo()
slide_niveles()
slide_motor()
slide_pantalla()
slide_valor()
slides_parametros()
slide_cierre()

prs.save(str(OUT))
print(f"OK · {len(prs.slides._sldIdLst)} slides (editable) -> {OUT}")
